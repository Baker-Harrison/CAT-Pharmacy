import tempfile
import unittest
from pathlib import Path

from backend import parser as parser_module

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - handled by skip
    Presentation = None


@unittest.skipIf(
    Presentation is None or not parser_module.HAS_PPTX,
    "python-pptx is required for parser tests",
)
class ParserTests(unittest.TestCase):
    def test_parse_pptx_extracts_units_and_objectives(self) -> None:
        presentation = Presentation()

        objectives_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        objectives_slide.shapes.title.text = "Learning Objectives"
        objectives_body = objectives_slide.shapes.placeholders[1].text_frame
        objectives_body.text = "Understand cardiac anatomy"
        objectives_body.add_paragraph().text = "Identify heart chambers"

        content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        content_slide.shapes.title.text = "Heart Anatomy"
        content_body = content_slide.shapes.placeholders[1].text_frame
        content_body.text = "Atria"
        content_body.add_paragraph().text = "Ventricles"

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "sample.pptx"
            presentation.save(pptx_path)

            units = parser_module.parse_pptx(str(pptx_path))

        self.assertEqual(len(units), 1)
        unit = units[0]
        self.assertEqual(unit.topic, "Heart Anatomy")
        self.assertIn("Atria", unit.key_points)
        self.assertIn("Ventricles", unit.key_points)
        self.assertIn("Understand cardiac anatomy", unit.learning_objectives)
        self.assertIn("Identify heart chambers", unit.learning_objectives)


if __name__ == "__main__":
    unittest.main()
