"""
Lecture Slide Parser for CAT-Pharmacy.
Extracts structured Knowledge Units from PowerPoint (.pptx) files using python-pptx.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import List

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    HAS_PPTX = True
except ImportError:  # pragma: no cover - guarded at runtime
    Presentation = None
    PP_PLACEHOLDER = None
    HAS_PPTX = False

from backend.models import KnowledgeUnit


NOISE_PATTERNS = [
    re.compile(r"(\d{3}-\d{3}-\d{4})", re.IGNORECASE),
    re.compile(r"Room\s+#?\d+", re.IGNORECASE),
    re.compile(r"@uams.edu", re.IGNORECASE),
    re.compile(r"www\.", re.IGNORECASE),
    re.compile(r"http", re.IGNORECASE),
]


def parse_pptx(file_path: str) -> List[KnowledgeUnit]:
    if not HAS_PPTX or Presentation is None:
        raise RuntimeError("python-pptx is required to parse PPTX files")
    presentation = Presentation(file_path)
    units: List[KnowledgeUnit] = []
    global_learning_objectives: List[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_content = _extract_slide_content(slide)

        if _is_learning_objectives_slide(slide_content.title):
            global_learning_objectives.extend(slide_content.key_points)
            continue

        if slide_content.title or slide_content.key_points:
            unit = KnowledgeUnit.create(
                topic=slide_content.title or f"Slide {slide_index}",
                subtopic="",
                source_slide_id=f"slide-{slide_index}",
                summary=slide_content.title or "",
                key_points=slide_content.key_points,
                learning_objectives=global_learning_objectives,
            )
            units.append(unit)

    return units


def _is_learning_objectives_slide(title: str | None) -> bool:
    if not title or not title.strip():
        return False

    lower_title = title.lower()
    return (
        "learning objective" in lower_title
        or "objectives" in lower_title
        or "learning outcome" in lower_title
        or "goals" in lower_title
    )


def _extract_slide_content(slide) -> "SlideContent":
    title = ""
    key_points: List[str] = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = _extract_text_from_shape(shape)
        if not text.strip():
            continue

        if _is_shape_title(shape) and not title:
            trimmed = text.strip()
            if not _is_noise(trimmed):
                title = trimmed
            continue

        for line in _split_lines(text):
            trimmed = line.strip()
            if trimmed and not _is_noise(trimmed):
                key_points.append(trimmed)

    return SlideContent(title=title, key_points=key_points)


def _is_shape_title(shape) -> bool:
    if not shape.is_placeholder:
        return False
    placeholder_type = shape.placeholder_format.type
    return placeholder_type in {
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
    }


def _extract_text_from_shape(shape) -> str:
    text_frame = shape.text_frame
    if text_frame is None:
        return ""

    lines: List[str] = []
    for paragraph in text_frame.paragraphs:
        runs_text = "".join(run.text or "" for run in paragraph.runs)
        lines.append(runs_text)
    return "\n".join(lines)


def _split_lines(text: str) -> List[str]:
    return [line for line in re.split(r"[\r\n]+", text) if line]


def _is_noise(text: str) -> bool:
    if not text or not text.strip():
        return True
    if len(text.strip()) < 2:
        return True
    return any(pattern.search(text) for pattern in NOISE_PATTERNS)


def _knowledge_unit_to_dict(unit: KnowledgeUnit) -> dict:
    return {
        "id": str(unit.id),
        "topic": unit.topic,
        "subtopic": unit.subtopic,
        "source_slide_id": unit.source_slide_id,
        "summary": unit.summary,
        "key_points": unit.key_points,
        "learning_objectives": unit.learning_objectives,
    }


class SlideContent:
    def __init__(self, title: str, key_points: List[str]) -> None:
        self.title = title
        self.key_points = key_points


def main() -> int:
    parser = argparse.ArgumentParser(description="Parse PPTX content into knowledge units.")
    parser.add_argument("file_path", help="Path to the .pptx file")
    args = parser.parse_args()

    file_path = Path(args.file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    units = parse_pptx(str(file_path))
    payload = [_knowledge_unit_to_dict(unit) for unit in units]
    print(json.dumps(payload, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
